import pdfplumber
from pptx import Presentation
import logging
from playwright.sync_api import sync_playwright
import trafilatura

logger = logging.getLogger(__name__)

import os

def load_pdf(path: str):
    pages = []
    filename = os.path.basename(path)
    with pdfplumber.open(path) as pdf:
        for i, page in enumerate(pdf.pages):
            # x_tolerance=1 helps prevent "C H A P T E R" issues by grouping closer characters
            text = page.extract_text(x_tolerance=1)
            if text:
                pages.append({
                    "source": f"{filename} (Page {i+1})",
                    "text": text
                })
    return pages


def load_ppt(path: str):
    prs = Presentation(path)
    slides = []
    filename = os.path.basename(path)

    for i, slide in enumerate(prs.slides):
        content = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                content.append(shape.text)

        if content:
            slides.append({
                "source": f"{filename} (Slide {i+1})",
                "text": " ".join(content)
            })

    return slides

def load_website(url: str):
    """
    Scrapes a website using Playwright for rendering and Trafilatura for extraction.
    Falls back to BeautifulSoup if Trafilatura fails.
    """
    from bs4 import BeautifulSoup
    import time
    
    logger.info(f"Scraping URL: {url}")
    print(f"[DEBUG] Starting scrape for: {url}")
    
    MAX_RETRIES = 3
    final_text = ""
    
    for attempt in range(MAX_RETRIES):
        try:
            with sync_playwright() as p:
                browser = p.chromium.launch(headless=True)
                context = browser.new_context(
                    user_agent="Mozilla/5.0 (Macintosh; Intel Mac OS X 10_15_7) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/120.0.0.0 Safari/537.36"
                )
                page = context.new_page()
                
                try:
                    page.goto(url, wait_until="domcontentloaded", timeout=45000) 
                    page.wait_for_load_state("domcontentloaded")
                    page.wait_for_timeout(2000) 
                except Exception as e:
                    print(f"[DEBUG] Page load timeout/error (attempt {attempt+1}/{MAX_RETRIES}): {e}")
                    if attempt < MAX_RETRIES - 1:
                        time.sleep(2)
                        continue

                html_content = page.content()
                browser.close()

                print(f"[DEBUG] Raw HTML captured. Length: {len(html_content)}")
                text = trafilatura.extract(html_content, include_comments=False, include_tables=True)
                
                if text and len(text) > 200:
                    print(f"[DEBUG] Trafilatura success. extracted {len(text)} chars.")
                    print(f"[DEBUG] Snippet: {text[:500]}...")
                    final_text = text
                    break
                
                print("[DEBUG] Trafilatura failed or too short. Falling back to BeautifulSoup.")

                soup = BeautifulSoup(html_content, "html.parser")
                
                for script in soup(["script", "style", "nav", "footer", "header", "aside"]):
                    script.decompose()


                target_tags = ['p', 'h1', 'h2', 'h3', 'h4', 'h5', 'h6', 'li', 'article', 'section']
                chunks = []
                for element in soup.find_all(target_tags):
                    content = element.get_text(strip=True)
                    if len(content) > 20: 
                        chunks.append(content)
                
                fallback_text = "\n\n".join(chunks)
                if fallback_text and len(fallback_text) > 200:
                     final_text = fallback_text
                     break
                
                print(f"[DEBUG] BeautifulSoup extraction too short. Retrying... ({attempt+1})")
                
        except Exception as e:
            print(f"[DEBUG] scraping error (attempt {attempt+1}): {e}")
            if attempt < MAX_RETRIES - 1:
                time.sleep(2)


    if not final_text or len(final_text) < 200:
        print(f"[DEBUG] Failed to extract meaningful content from {url}. Final text length: {len(final_text) if final_text else 0}")
        return [] 
    

    block_phrases = ["please verify you are a human", "access denied", "captcha"]
    if any(phrase in final_text.lower() for phrase in block_phrases):
         print(f"[DEBUG] Detected bot block/captcha. Phrases found: {[p for p in block_phrases if p in final_text.lower()]}")
         return []
         
    # "Enable javascript" is common in footnotes/noscript tags of valid sites. Only block if text is short.
    if "enable javascript" in final_text.lower() and len(final_text) < 1000:
        print("[DEBUG] Detected 'enable javascript' in short text. Treating as bot block.")
        return []

    return [{"source": url, "text": final_text}]
            


